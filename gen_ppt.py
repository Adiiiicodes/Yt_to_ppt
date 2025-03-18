from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_beautiful_pptx(transcription):
    # Create a presentation object
    prs = Presentation()
    # Slide Layout (0 - Title Slide, 1 - Title and Content)
    slide_layout = prs.slide_layouts[1]
    # Add a title slide with a background color
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Set background color for title slide with gradient effect
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[0].color.rgb = RGBColor(0, 102, 204)  # Deeper blue
    fill.gradient_stops[1].position = 1
    fill.gradient_stops[1].color.rgb = RGBColor(135, 206, 250)  # Light blue
    
    # Add title text AFTER setting background
    title = slide.shapes.title
    title.text = "Transcription Overview"
    subtitle = slide.placeholders[1]
    subtitle.text = "Generated from YouTube transcription"
    
    # Make sure text is visible against blue background
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        paragraph.font.bold = True
    
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
    
    # Split transcription into chunks of 70 words each
    words = transcription.split()
    chunk_size = 70  # We want 70 words per slide (2 bullet points of 70 words)
    chunks = [words[i:i + chunk_size] for i in range(0, len(words), chunk_size)]
    
    # Define gradient color schemes for content slides - using very light gradients for readability
    gradient_schemes = [
        # Light blue gradient
        (RGBColor(240, 248, 255), RGBColor(220, 235, 250)),
        # Light green gradient
        (RGBColor(240, 255, 240), RGBColor(220, 250, 230)),
        # Light purple gradient
        (RGBColor(250, 240, 255), RGBColor(235, 225, 250)),
        # Light gold gradient
        (RGBColor(255, 250, 240), RGBColor(250, 240, 220))
    ]
    
    # Process each chunk for a new slide
    for idx, chunk in enumerate(chunks):
        # Create a new slide for each chunk of transcription
        slide = prs.slides.add_slide(slide_layout)
        
        # Set gradient background for content slides
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].position = 0
        fill.gradient_stops[0].color.rgb = gradient_schemes[idx % len(gradient_schemes)][0]
        fill.gradient_stops[1].position = 1
        fill.gradient_stops[1].color.rgb = gradient_schemes[idx % len(gradient_schemes)][1]
        
        # Add a slide title AFTER setting background
        title = slide.shapes.title
        title.text = f"Video Insights part {idx + 1}"
        
        # Add accent color to title
        accent_colors = [
            RGBColor(0, 102, 204),  # Blue
            RGBColor(0, 128, 64),   # Green
            RGBColor(128, 0, 128),  # Purple
            RGBColor(184, 134, 11)  # Dark gold
        ]
        
        # Apply accent color to title text
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.color.rgb = accent_colors[idx % len(accent_colors)]
            paragraph.font.bold = True
        
        # Add body content (split into two points of 70 words each)
        body_text = slide.placeholders[1]
        text_frame = body_text.text_frame
        text_frame.clear()  # Clear any default text
        
        # Add two bullet points for each chunk
        # First Bullet Point
        paragraph1 = text_frame.add_paragraph()
        paragraph1.text = ' '.join(chunk[:35]) if len(chunk) > 35 else ' '.join(chunk)
        paragraph1.level = 0  # Top level bullet
        
        # Second Bullet Point (only if we have enough words)
        if len(chunk) > 35:
            paragraph2 = text_frame.add_paragraph()
            paragraph2.text = ' '.join(chunk[35:])
            paragraph2.level = 0  # Top level bullet
        
        # Text formatting for both paragraphs - dark text for readability
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)  # Change font size
            paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text for maximum readability
            paragraph.font.bold = False  # Not bold
    
    # Save the presentation
    pptx_filename = "enhanced_presentation.pptx"
    prs.save(pptx_filename)
    print(f"Presentation saved as {pptx_filename}")
